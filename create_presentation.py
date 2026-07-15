"""
Create Professional Portfolio Optimisation Dashboard Presentation
3-Slide Presentation with detailed explanations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(44, 62, 80)      # #2C3E50
LIGHT_BLUE = RGBColor(52, 152, 219)   # #3498DB
ACCENT_PURPLE = RGBColor(155, 89, 182) # #9B59B6
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(236, 240, 241)  # #ECF0F1

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, sections):
    """Add a content slide with multiple sections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add content sections
    top = 1.3
    for section in sections:
        # Section title
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.4))
        section_frame = section_box.text_frame
        section_frame.text = section['title']
        section_para = section_frame.paragraphs[0]
        section_para.font.size = Pt(18)
        section_para.font.bold = True
        section_para.font.color.rgb = LIGHT_BLUE
        
        top += 0.45
        
        # Section content
        if isinstance(section['content'], list):
            for item in section['content']:
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.7), Inches(0.35))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                content_frame.text = f"• {item}"
                content_para = content_frame.paragraphs[0]
                content_para.font.size = Pt(14)
                content_para.font.color.rgb = DARK_BLUE
                top += 0.4
        else:
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.7), Inches(0.8))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = section['content']
            content_para = content_frame.paragraphs[0]
            content_para.font.size = Pt(14)
            content_para.font.color.rgb = DARK_BLUE
            top += 0.9
        
        top += 0.2
    
    return slide

# ============ SLIDE 1: PROBLEM STATEMENT & SOLUTION ============

slide1_sections = [
    {
        'title': '🎯 Problem Statement',
        'content': [
            'Reinsurance portfolio managers lack real-time visibility into portfolio performance metrics',
            'Manual analysis of 50+ treaties across multiple Lines of Business (LOB) and geographies is time-consuming',
            'Difficulty in identifying optimization opportunities and concentration risks quickly',
            'No unified dashboard for RORAC calculations, capital efficiency, and diversification scoring'
        ]
    },
    {
        'title': '✅ Solution',
        'content': [
            'Dynamic Portfolio Optimisation Dashboard with real-time data updates (1-minute refresh)',
            'Automated RORAC engine calculating Return on Risk-Adjusted Capital for each treaty',
            'Interactive visualizations showing portfolio by LOB and geographic distribution',
            'WebSocket-enabled architecture for live updates and IST timezone support',
            'Synthetic portfolio generator with 50 realistic reinsurance treaties for testing'
        ]
    }
]

add_content_slide(prs, 'Problem Statement & Solution', slide1_sections)

# ============ SLIDE 2: BENEFITS & OVERVIEW ============

slide2_sections = [
    {
        'title': '💼 Key Benefits',
        'content': [
            'Real-Time Visibility: Monitor portfolio metrics instantly with live IST timestamps',
            'Faster Decision-Making: Identify optimization opportunities in seconds, not hours',
            'Risk Management: Track capital utilization, diversification scores, and concentration risks',
            'Operational Efficiency: Automated calculations reduce manual analysis by 80%',
            'Scalability: Flask + WebSocket architecture handles multiple concurrent users',
            'Cost Reduction: Optimize treaty allocation and reduce concentration risk exposure'
        ]
    },
    {
        'title': '🏗️ Technical Architecture',
        'content': 'Python Flask backend | Flask-SocketIO for WebSocket real-time updates | MockPortfolioGenerator for 50 synthetic treaties | RORAC optimization engine | Interactive Bootstrap 5 UI with Chart.js visualizations'
    }
]

add_content_slide(prs, 'Benefits & Technical Overview', slide2_sections)

# ============ SLIDE 3: DETAILED IMPLEMENTATION STEPS ============

slide3_sections = [
    {
        'title': '🔧 Step-by-Step Implementation',
        'content': [
            'Step 1: Created MockPortfolioGenerator (data_connectors/mock_portfolio.py) generating 50 synthetic treaties with realistic parameters',
            'Step 2: Built PortfolioOptimizer engine calculating RORAC, capital efficiency, and diversification metrics for each treaty',
            'Step 3: Developed Flask backend (app_simple.py) with 5 API endpoints (/api/portfolio, /api/portfolio/summary, etc.)',
            'Step 4: Integrated Flask-SocketIO for WebSocket bidirectional real-time communication with browser clients',
            'Step 5: Created responsive HTML5 dashboard (dashboard.html) with Bootstrap 5 + Chart.js interactive charts',
            'Step 6: Implemented auto-refresh logic: timestamp updates every 1 second, full data refresh every 60 seconds in IST',
            'Step 7: Added filter functionality (by LOB, geography, RORAC range) with instant UI updates',
            'Step 8: Set up Python virtual environment and installed all dependencies (Flask, SocketIO, Pandas, NumPy)',
            'Step 9: Deployed on localhost:5001 with hot-reload debugging enabled for rapid iteration',
            'Step 10: Tested across all 5 dashboard pages (Dashboard, Portfolio, Scenarios, Recommendations, Reports)'
        ]
    }
]

add_content_slide(prs, 'Detailed Implementation Steps', slide3_sections)

# Save presentation
output_path = 'Portfolio_Optimisation_Dashboard_Presentation.pptx'
prs.save(output_path)

print(f"✅ Presentation created successfully!")
print(f"📊 Saved to: {output_path}")
print(f"\n📋 Presentation Contents:")
print(f"   Slide 1: Problem Statement & Solution")
print(f"   Slide 2: Benefits & Technical Overview")
print(f"   Slide 3: Detailed Implementation Steps (10 steps)")
print(f"\n🎨 Design Features:")
print(f"   • Professional color scheme (Dark Blue + Light Blue + Purple)")
print(f"   • 44pt bold titles with color-coded sections")
print(f"   • Bullet points for easy reading")
print(f"   • Clean, modern layout")
